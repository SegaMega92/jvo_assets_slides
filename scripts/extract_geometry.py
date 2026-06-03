"""Шаг 7: позиции и размеры всех shapes на каждом слайде.

EMU → px при 96dpi: emu / 914400 * 96
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SOURCE_DIR = Path(__file__).parent.parent / "source_pptx"
OUT_FILE = Path(__file__).parent.parent / "analysis" / "geometry_dump.json"


def emu_to_px(emu: int) -> float:
    return emu / 914400 * 96


def shape_type_name(shape) -> str:
    try:
        return str(shape.shape_type).split(".")[-1] if shape.shape_type else "UNKNOWN"
    except Exception:
        return "UNKNOWN"


def placeholder_info(shape) -> dict | None:
    try:
        if not shape.is_placeholder:
            return None
        ph = shape.placeholder_format
        return {"type": str(ph.type) if ph.type else None, "idx": ph.idx}
    except Exception:
        return None


def text_preview(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text[:80].replace("\n", " | ")


def walk(shapes, slide_idx, file_name, slide_w_px, slide_h_px, records, depth=0):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk(shape.shapes, slide_idx, file_name, slide_w_px, slide_h_px,
                 records, depth + 1)
            continue

        try:
            left_px = emu_to_px(shape.left or 0)
            top_px = emu_to_px(shape.top or 0)
            w_px = emu_to_px(shape.width or 0)
            h_px = emu_to_px(shape.height or 0)
        except Exception:
            continue

        records.append({
            "file": file_name,
            "slide": slide_idx,
            "shape": shape.name,
            "shape_type": shape_type_name(shape),
            "placeholder": placeholder_info(shape),
            "depth": depth,
            "left_px": round(left_px, 1),
            "top_px": round(top_px, 1),
            "w_px": round(w_px, 1),
            "h_px": round(h_px, 1),
            "left_pct": round(left_px / slide_w_px * 100, 1),
            "top_pct": round(top_px / slide_h_px * 100, 1),
            "w_pct": round(w_px / slide_w_px * 100, 1),
            "h_pct": round(h_px / slide_h_px * 100, 1),
            "text_preview": text_preview(shape),
        })


def main() -> None:
    records = []
    for f in sorted(SOURCE_DIR.glob("*.pptx")):
        prs = Presentation(str(f))
        sw_px = emu_to_px(prs.slide_width)
        sh_px = emu_to_px(prs.slide_height)
        for i, slide in enumerate(prs.slides):
            walk(slide.shapes, i, f.name, sw_px, sh_px, records)

    OUT_FILE.write_text(json.dumps(records, indent=2, ensure_ascii=False))

    placeholders = [r for r in records if r["placeholder"]]
    print(f"Total shapes: {len(records)}")
    print(f"Placeholders: {len(placeholders)}")
    from collections import Counter
    ph_types = Counter(r["placeholder"]["type"] for r in placeholders)
    print(f"Placeholder types: {dict(ph_types)}")
    shape_types = Counter(r["shape_type"] for r in records)
    print(f"Top shape types: {dict(shape_types.most_common(8))}")
    print(f"→ {OUT_FILE}")


if __name__ == "__main__":
    main()
