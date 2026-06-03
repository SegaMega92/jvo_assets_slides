"""Шаг 4: извлечь все hex-цвета со всех shapes на всех слайдах.

Источники цвета:
  - shape.fill (solid fill)
  - shape.line.color (border)
  - text run font.color (текст)
  - slide background

Theme colors разрешаются через theme_colors.json (шаг 3).
"""
import json
from pathlib import Path
from xml.etree import ElementTree as ET
from pptx import Presentation
from pptx.dml.color import RGBColor

SOURCE_DIR = Path(__file__).parent.parent / "source_pptx"
ANALYSIS_DIR = Path(__file__).parent.parent / "analysis"
THEME_FILE = ANALYSIS_DIR / "theme_colors.json"
OUT_FILE = ANALYSIS_DIR / "all_colors.json"

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

THEME_MAP = {
    "bg1": "lt1", "tx1": "dk1", "bg2": "lt2", "tx2": "dk2",
    "accent1": "accent1", "accent2": "accent2", "accent3": "accent3",
    "accent4": "accent4", "accent5": "accent5", "accent6": "accent6",
    "hlink": "hlink", "folHlink": "folHlink",
    "dk1": "dk1", "lt1": "lt1", "dk2": "dk2", "lt2": "lt2",
}


def apply_lum_mod(hex_color: str, lum_mod: float | None, lum_off: float | None) -> str:
    """Применить lumMod/lumOff к hex цвету (приближённо).

    PowerPoint считает в HSL. Здесь упрощённо: lumMod умножает яркость,
    lumOff прибавляет. Для точности нужен HSL-конверт, но для статистики хватит.
    """
    if not hex_color or (lum_mod is None and lum_off is None):
        return hex_color
    r = int(hex_color[1:3], 16) / 255
    g = int(hex_color[3:5], 16) / 255
    b = int(hex_color[5:7], 16) / 255
    if lum_mod is not None:
        r *= lum_mod
        g *= lum_mod
        b *= lum_mod
    if lum_off is not None:
        r += lum_off
        g += lum_off
        b += lum_off
    r = max(0, min(1, r))
    g = max(0, min(1, g))
    b = max(0, min(1, b))
    return f"#{int(r*255):02X}{int(g*255):02X}{int(b*255):02X}"


def resolve_color_elem(color_elem, theme: dict) -> str | None:
    """Разобрать <a:solidFill>, <a:srgbClr>, <a:schemeClr>... → hex."""
    if color_elem is None:
        return None

    # srgbClr — прямой hex
    srgb = color_elem.find(".//a:srgbClr", NS)
    if srgb is not None:
        hex_ = "#" + srgb.get("val").upper()
        lm = srgb.find("a:lumMod", NS)
        lo = srgb.find("a:lumOff", NS)
        lm_v = int(lm.get("val")) / 100000 if lm is not None else None
        lo_v = int(lo.get("val")) / 100000 if lo is not None else None
        return apply_lum_mod(hex_, lm_v, lo_v)

    # schemeClr — ссылка на тему
    sch = color_elem.find(".//a:schemeClr", NS)
    if sch is not None:
        val = sch.get("val")
        mapped = THEME_MAP.get(val, val)
        if mapped in theme:
            hex_ = theme[mapped]
            lm = sch.find("a:lumMod", NS)
            lo = sch.find("a:lumOff", NS)
            lm_v = int(lm.get("val")) / 100000 if lm is not None else None
            lo_v = int(lo.get("val")) / 100000 if lo is not None else None
            return apply_lum_mod(hex_, lm_v, lo_v)

    # sysClr
    sys = color_elem.find(".//a:sysClr", NS)
    if sys is not None:
        last = sys.get("lastClr")
        if last:
            return "#" + last.upper()

    return None


def extract_shape_fill(shape, theme: dict) -> str | None:
    """Solid fill из shape."""
    try:
        sp_pr = shape.element.find(".//p:spPr", NS)
        if sp_pr is None:
            return None
        solid = sp_pr.find("a:solidFill", NS)
        return resolve_color_elem(solid, theme) if solid is not None else None
    except Exception:
        return None


def extract_shape_line(shape, theme: dict) -> str | None:
    try:
        sp_pr = shape.element.find(".//p:spPr", NS)
        if sp_pr is None:
            return None
        ln = sp_pr.find("a:ln", NS)
        if ln is None:
            return None
        solid = ln.find("a:solidFill", NS)
        return resolve_color_elem(solid, theme) if solid is not None else None
    except Exception:
        return None


def extract_run_color(run, theme: dict) -> str | None:
    """Цвет шрифта из run."""
    try:
        rpr = run._r.find("a:rPr", NS)
        if rpr is None:
            return None
        solid = rpr.find("a:solidFill", NS)
        return resolve_color_elem(solid, theme) if solid is not None else None
    except Exception:
        return None


def walk_shapes(shapes, slide_idx, file_name, theme, records):
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            walk_shapes(shape.shapes, slide_idx, file_name, theme, records)
            continue

        fill = extract_shape_fill(shape, theme)
        if fill:
            records.append({
                "file": file_name, "slide": slide_idx,
                "shape": shape.name, "role": "fill", "hex": fill,
            })

        line = extract_shape_line(shape, theme)
        if line:
            records.append({
                "file": file_name, "slide": slide_idx,
                "shape": shape.name, "role": "line", "hex": line,
            })

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    color = extract_run_color(run, theme)
                    if color:
                        records.append({
                            "file": file_name, "slide": slide_idx,
                            "shape": shape.name, "role": "text", "hex": color,
                            "text_preview": run.text[:30],
                        })


def extract_slide_bg(slide, theme: dict) -> str | None:
    """Цвет фона слайда."""
    bg = slide.element.find(".//p:bg", NS)
    if bg is None:
        return None
    solid = bg.find(".//a:solidFill", NS)
    return resolve_color_elem(solid, theme) if solid is not None else None


def main() -> None:
    themes = {t["file"]: t["colors"] for t in json.loads(THEME_FILE.read_text())}
    records = []

    for f in sorted(SOURCE_DIR.glob("*.pptx")):
        prs = Presentation(str(f))
        theme = themes.get(f.name, {})
        for i, slide in enumerate(prs.slides):
            bg = extract_slide_bg(slide, theme)
            if bg:
                records.append({
                    "file": f.name, "slide": i,
                    "shape": "_slide_bg", "role": "fill", "hex": bg,
                })
            walk_shapes(slide.shapes, i, f.name, theme, records)

    # Подсчёт уникальных hex
    hexes = [r["hex"] for r in records]
    unique = sorted(set(hexes))
    print(f"Total color occurrences: {len(records)}")
    print(f"Unique hex values: {len(unique)}")
    print(f"By role: fill={sum(1 for r in records if r['role']=='fill')}, "
          f"line={sum(1 for r in records if r['role']=='line')}, "
          f"text={sum(1 for r in records if r['role']=='text')}")

    OUT_FILE.write_text(json.dumps(records, indent=2, ensure_ascii=False))
    print(f"→ {OUT_FILE}")


if __name__ == "__main__":
    main()
