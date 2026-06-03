"""Шаг 6: извлечь типографику со всех runs.

Иерархия наследования:
  run.font → paragraph (через listStyle) → text_frame → placeholder → layout → master

python-pptx даёт нам run.font, но если поле None — нужно
смотреть в layout и master placeholder с тем же idx/type.
"""
import json
from pathlib import Path
from xml.etree import ElementTree as ET
from pptx import Presentation
from pptx.util import Pt

SOURCE_DIR = Path(__file__).parent.parent / "source_pptx"
OUT_FILE = Path(__file__).parent.parent / "analysis" / "typography_dump.json"

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def get_font_attrs_from_rpr(rpr) -> dict:
    """Достать font name/size/bold/italic/color из <a:rPr>."""
    if rpr is None:
        return {}
    out = {}
    # size in 1/100 pt
    sz = rpr.get("sz")
    if sz:
        out["size_pt"] = int(sz) / 100
    b = rpr.get("b")
    if b is not None:
        out["bold"] = b == "1"
    i = rpr.get("i")
    if i is not None:
        out["italic"] = i == "1"
    # font family
    latin = rpr.find("a:latin", NS)
    if latin is not None:
        out["font"] = latin.get("typeface")
    # spacing (трекинг)
    spc = rpr.get("spc")
    if spc is not None:
        out["letter_spacing_pt"] = int(spc) / 100
    return out


def get_lvl_props(text_frame, lvl: int) -> dict:
    """Из text_frame.lstStyle взять уровневые свойства."""
    try:
        body_pr = text_frame._txBody.find("a:lstStyle", NS)
        if body_pr is None:
            return {}
        lvl_pPr = body_pr.find(f"a:lvl{lvl+1}pPr", NS) if lvl > 0 else body_pr.find("a:lvl1pPr", NS)
        if lvl_pPr is None:
            return {}
        def_rpr = lvl_pPr.find("a:defRPr", NS)
        return get_font_attrs_from_rpr(def_rpr) if def_rpr is not None else {}
    except Exception:
        return {}


def get_placeholder_info(shape) -> dict | None:
    """Если shape — placeholder, вернуть type+idx."""
    try:
        if not shape.is_placeholder:
            return None
        ph = shape.placeholder_format
        return {
            "type": str(ph.type) if ph.type else None,
            "idx": ph.idx,
        }
    except Exception:
        return None


def get_master_layout_defaults(prs) -> dict:
    """Из master выдернуть дефолты для title/body/other."""
    out = {}
    master_elem = prs.slide_master.element
    txStyles = master_elem.find(".//p:txStyles", NS)
    if txStyles is None:
        return out

    for tag, key in [("p:titleStyle", "title"), ("p:bodyStyle", "body"), ("p:otherStyle", "other")]:
        style_elem = txStyles.find(tag, NS)
        if style_elem is None:
            continue
        levels = {}
        for lvl in range(1, 10):
            lvl_pPr = style_elem.find(f"a:lvl{lvl}pPr", NS)
            if lvl_pPr is None:
                continue
            def_rpr = lvl_pPr.find("a:defRPr", NS)
            attrs = get_font_attrs_from_rpr(def_rpr)
            if attrs:
                levels[f"lvl{lvl}"] = attrs
        if levels:
            out[key] = levels
    return out


def walk_runs(shapes, slide_idx, file_name, records, master_defaults):
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            walk_runs(shape.shapes, slide_idx, file_name, records, master_defaults)
            continue
        if not shape.has_text_frame:
            continue

        ph_info = get_placeholder_info(shape)
        tf = shape.text_frame

        for p_idx, para in enumerate(tf.paragraphs):
            lvl = para.level
            for r_idx, run in enumerate(para.runs):
                text = run.text
                if not text.strip():
                    continue

                rpr = run._r.find("a:rPr", NS)
                attrs = get_font_attrs_from_rpr(rpr)

                # цвет отдельно (вложенный solidFill)
                color = None
                if rpr is not None:
                    solid = rpr.find("a:solidFill", NS)
                    if solid is not None:
                        srgb = solid.find("a:srgbClr", NS)
                        if srgb is not None:
                            color = "#" + srgb.get("val").upper()

                # дефолт по уровню из text_frame
                if "size_pt" not in attrs:
                    lvl_attrs = get_lvl_props(tf, lvl)
                    if "size_pt" in lvl_attrs:
                        attrs["size_pt"] = lvl_attrs["size_pt"]
                    if "font" not in attrs and "font" in lvl_attrs:
                        attrs["font"] = lvl_attrs["font"]

                # дефолт из master по типу placeholder
                if ph_info and ("size_pt" not in attrs or "font" not in attrs):
                    is_title = ph_info["type"] and "TITLE" in ph_info["type"]
                    master_key = "title" if is_title else "body"
                    master_lvls = master_defaults.get(master_key, {})
                    master_attrs = master_lvls.get(f"lvl{lvl+1}", {})
                    if "size_pt" not in attrs and "size_pt" in master_attrs:
                        attrs["size_pt"] = master_attrs["size_pt"]
                        attrs["_inherited_size"] = True
                    if "font" not in attrs and "font" in master_attrs:
                        attrs["font"] = master_attrs["font"]
                        attrs["_inherited_font"] = True

                records.append({
                    "file": file_name,
                    "slide": slide_idx,
                    "shape": shape.name,
                    "placeholder": ph_info,
                    "level": lvl,
                    "text_preview": text[:60],
                    "color": color,
                    **attrs,
                })


def main() -> None:
    records = []
    for f in sorted(SOURCE_DIR.glob("*.pptx")):
        prs = Presentation(str(f))
        master_defaults = get_master_layout_defaults(prs)
        for i, slide in enumerate(prs.slides):
            walk_runs(slide.shapes, i, f.name, records, master_defaults)

    OUT_FILE.write_text(json.dumps(records, indent=2, ensure_ascii=False))

    sizes = [r.get("size_pt") for r in records if r.get("size_pt")]
    fonts = [r.get("font") for r in records if r.get("font")]
    print(f"Records: {len(records)}")
    print(f"With size: {len(sizes)}, unique sizes: {sorted(set(sizes))}")
    from collections import Counter
    print(f"Fonts: {Counter(fonts).most_common()}")
    print(f"→ {OUT_FILE}")


if __name__ == "__main__":
    main()
