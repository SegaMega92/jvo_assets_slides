"""Шаг 8: разобрать slide master и slide layouts.

Master = базовая страница со всеми дефолтами.
Layouts = варианты разметки (Title Slide, Title + Content, ...).
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SOURCE_DIR = Path(__file__).parent.parent / "source_pptx"
OUT_FILE = Path(__file__).parent.parent / "analysis" / "master_layouts.json"


def emu_to_px(emu: int) -> float:
    return (emu or 0) / 914400 * 96


def shape_summary(shape, sw_px: float, sh_px: float, depth: int = 0) -> dict:
    out = {
        "name": shape.name,
        "shape_type": str(shape.shape_type).split(".")[-1] if shape.shape_type else None,
        "depth": depth,
    }
    try:
        out["left_pct"] = round(emu_to_px(shape.left) / sw_px * 100, 1)
        out["top_pct"] = round(emu_to_px(shape.top) / sh_px * 100, 1)
        out["w_pct"] = round(emu_to_px(shape.width) / sw_px * 100, 1)
        out["h_pct"] = round(emu_to_px(shape.height) / sh_px * 100, 1)
    except Exception:
        pass

    try:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            out["placeholder"] = {
                "type": str(ph.type) if ph.type else None,
                "idx": ph.idx,
            }
    except Exception:
        pass

    if shape.has_text_frame and shape.text_frame.text.strip():
        out["text_preview"] = shape.text_frame.text[:60].replace("\n", " | ")

    return out


def collect_shapes(shapes, sw_px: float, sh_px: float, depth: int = 0) -> list:
    out = []
    for shape in shapes:
        out.append(shape_summary(shape, sw_px, sh_px, depth))
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            out.extend(collect_shapes(shape.shapes, sw_px, sh_px, depth + 1))
    return out


def layout_usage_count(prs) -> dict:
    """Сколько слайдов использует каждый layout."""
    counts = {}
    for slide in prs.slides:
        name = slide.slide_layout.name
        counts[name] = counts.get(name, 0) + 1
    return counts


def main() -> None:
    results = []
    for f in sorted(SOURCE_DIR.glob("*.pptx")):
        prs = Presentation(str(f))
        sw_px = emu_to_px(prs.slide_width)
        sh_px = emu_to_px(prs.slide_height)

        master = prs.slide_master
        master_shapes = collect_shapes(master.shapes, sw_px, sh_px)

        usage = layout_usage_count(prs)

        layouts_info = []
        for layout in master.slide_layouts:
            layouts_info.append({
                "name": layout.name,
                "used_by_slides": usage.get(layout.name, 0),
                "shapes": collect_shapes(layout.shapes, sw_px, sh_px),
            })

        results.append({
            "file": f.name,
            "slide_w_px": round(sw_px),
            "slide_h_px": round(sh_px),
            "master_shapes": master_shapes,
            "layouts": layouts_info,
        })

    OUT_FILE.write_text(json.dumps(results, indent=2, ensure_ascii=False))

    print(f"\n=== Master / Layout summary ===\n")
    for r in results:
        print(f"\n📄 {r['file']}")
        print(f"   slide: {r['slide_w_px']}×{r['slide_h_px']}px")
        print(f"   master shapes: {len(r['master_shapes'])}")
        print(f"   layouts ({len(r['layouts'])}):")
        for l in r["layouts"]:
            mark = "★" if l["used_by_slides"] > 0 else " "
            print(f"     {mark} {l['name']:<45} used: {l['used_by_slides']:2d}  shapes: {len(l['shapes'])}")

    print(f"\n→ {OUT_FILE}")


if __name__ == "__main__":
    main()
